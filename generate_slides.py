import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Function to get the chart type object from a string
def get_chart_type(chart_type_str):
    chart_types = {
        'bar-horizontal': XL_CHART_TYPE.BAR_CLUSTERED,
        'bar-vertical': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        # Add any other chart types you need
    }
    return chart_types.get(chart_type_str.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)


# Function to find slide layout by name
def find_layout_by_name(prs, layout_name):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    return None

# Load the dataset from Excel
excel_file = 'data.xlsx'
df = pd.read_excel(excel_file)

# Load the presentation template
template_pptx = 'template.pptx'
prs = Presentation(template_pptx)

# Process each unique question_id
for question_id in df['question_id'].unique():
    question_data = df[df['question_id'] == question_id]
    question_text = question_data['question_text'].iloc[0]
    text_summary = question_data['text_summary'].iloc[0]  # Get the text_summary
    chart_type_str = question_data['chart_type'].iloc[0]
    chart_type = get_chart_type(chart_type_str)
    chart_layout_name = question_data['chart_layout'].iloc[0]

    # Find the slide layout by name and add a new slide
    slide_layout = find_layout_by_name(prs, chart_layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # Set the title for the slide
    slide.shapes.title.text = question_text

    # Add chart
    chart_data = CategoryChartData()
    chart_data.categories = question_data['response'].tolist()
    chart_data.add_series('Series 1', (question_data['value']*100).tolist())  # Multiply by 100 to convert to percentage

    # Insert text_summary into the specified text placeholder
    text_placeholder = slide.placeholders[11]  # Index 11 for the text_summary placeholder
    text_placeholder.text = text_summary  # Set the text_summary text

    # Choose the chart placeholder by index and insert the chart
    chart_placeholder = slide.placeholders[10]  # This is the placeholder index for the chart
    chart_frame = chart_placeholder.insert_chart(chart_type, chart_data)
    chart = chart_frame.chart

   # Formatting based on the chart type
    if chart_type in (XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.LINE_MARKERS):
        for series in chart.series:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x14, 0x60, 0x82)  # Default blue color

            # Add data labels formatted as percentages with no decimal points
            series.has_data_labels = True
            for point in series.points:
                point.data_label.number_format = '0%'
                point.data_label.font.size = Pt(10)
                point.data_label.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black font color
                
            # Remove the chart title and legend
            chart.has_title = False
            chart.has_legend = False

            # Remove the chart axis
            chart.value_axis.visible = False
            chart.category_axis.visible = True

        # Remove gridlines for all charts
        if chart.category_axis and chart.category_axis.has_major_gridlines:
            chart.category_axis.major_gridlines.format.line.fill.background()

        if chart.value_axis and chart.value_axis.has_major_gridlines:
            chart.value_axis.major_gridlines.format.line.fill.background()


# Save the presentation
prs.save('output_presentation.pptx')

print("Presentation created successfully!")
